from __future__ import annotations

import re
import shutil
import uuid
from pathlib import Path

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent.parent
STATIC_DIR = BASE_DIR / "static"
TEMPLATE_DIR = BASE_DIR / "storage" / "templates"
GENERATED_DIR = BASE_DIR / "storage" / "generated"

for directory in (TEMPLATE_DIR, GENERATED_DIR):
    directory.mkdir(parents=True, exist_ok=True)

app = FastAPI(title="GenerateX", version="0.1.0")


def safe_filename(value: str) -> str:
    """Create a safe filename while retaining a readable participant name."""
    cleaned = re.sub(r"[<>:\\"/\\|?*]+", "", value).strip()
    cleaned = re.sub(r"\s+", "_", cleaned)
    return cleaned[:100] or "certificate"


def replace_placeholder_in_shape(shape, placeholder: str, replacement: str) -> int:
    """Replace a placeholder in text frames, including grouped shapes."""
    count = 0
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            count += replace_placeholder_in_shape(child, placeholder, replacement)
        return count

    if not shape.has_text_frame:
        return count

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if placeholder in run.text:
                count += run.text.count(placeholder)
                run.text = run.text.replace(placeholder, replacement)
    return count


def fill_presentation(template_path: Path, output_path: Path, name: str) -> int:
    presentation = Presentation(template_path)
    replacements = 0

    for slide in presentation.slides:
        for shape in slide.shapes:
            replacements += replace_placeholder_in_shape(shape, "{{NAME}}", name)
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "{{NAME}}" in run.text:
                                replacements += run.text.count("{{NAME}}")
                                run.text = run.text.replace("{{NAME}}", name)

    presentation.save(output_path)
    return replacements


@app.get("/api/health")
def health_check():
    return {"status": "ok"}


@app.post("/api/generate")
async def generate_certificate(
    name: str = Form(...), template: UploadFile = File(...)
):
    clean_name = name.strip()
    if not clean_name:
        raise HTTPException(status_code=400, detail="Please enter a participant name.")
    if not template.filename or not template.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Please upload a PowerPoint (.pptx) template.")

    request_id = uuid.uuid4().hex
    template_path = TEMPLATE_DIR / f"{request_id}.pptx"
    output_path = GENERATED_DIR / f"{safe_filename(clean_name)}_Certificate.pptx"

    try:
        with template_path.open("wb") as target:
            shutil.copyfileobj(template.file, target)
        replacements = fill_presentation(template_path, output_path, clean_name)
    except Exception as exc:
        raise HTTPException(
            status_code=400,
            detail="We could not read that PowerPoint file. Please upload a valid .pptx template.",
        ) from exc
    finally:
        await template.close()
        template_path.unlink(missing_ok=True)

    headers = {"X-GenerateX-Replacements": str(replacements)}
    return FileResponse(
        output_path,
        media_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        filename=output_path.name,
        headers=headers,
    )


app.mount("/", StaticFiles(directory=STATIC_DIR, html=True), name="static")
